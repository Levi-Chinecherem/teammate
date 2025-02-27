from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Q1 Overview"
subtitle = slide.placeholders[1]
subtitle.text = "Sales up 10%"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Next Steps"
content = slide.shapes.placeholders[1]
content.text = "Expand marketing"

prs.save("../data/q1_slides.pptx")
print("Sample PPTX created at data/q1_slides.pptx")