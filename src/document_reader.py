import os
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image
import speech_recognition as sr
from moviepy.editor import VideoFileClip
from transformers import pipeline
from utils.db import get_db_connection
from msgraph.core import GraphClient
import requests
from dotenv import load_dotenv
import logging

# Load environment variables
load_dotenv()
client_id = os.getenv("CLIENT_ID")
client_secret = os.getenv("CLIENT_SECRET")
tenant_id = os.getenv("TENANT_ID")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("DocumentReader")

# Graph API client
def get_graph_client():
    auth_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    auth_data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": "https://graph.microsoft.com/.default"
    }
    response = requests.post(auth_url, data=auth_data)
    token = response.json().get("access_token")
    return GraphClient(api_version="v1.0", access_token=token)

# NLP for question answering
qa_pipeline = pipeline("question-answering", model="distilbert-base-uncased-distilled-squad")

def read_document(file_path: str) -> tuple[str, str]:
    """Read content from various document types."""
    file_type = file_path.split(".")[-1].lower()
    content = ""

    try:
        if file_type == "xlsx":
            df = pd.read_excel(file_path)
            content = df.to_string()
        elif file_type == "csv":
            df = pd.read_csv(file_path)
            content = df.to_string()
        elif file_type == "docx":
            doc = Document(file_path)
            content = " ".join([para.text for para in doc.paragraphs])
        elif file_type == "txt":
            with open(file_path, "r") as f:
                content = f.read()
        elif file_type == "pdf":
            reader = PdfReader(file_path)
            content = " ".join([page.extract_text() for page in reader.pages])
        elif file_type == "pptx":
            prs = Presentation(file_path)
            content = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])
        elif file_type in ["jpg", "png"]:
            content = pytesseract.image_to_string(Image.open(file_path))
        elif file_type in ["mp3", "wav"]:
            r = sr.Recognizer()
            with sr.AudioFile(file_path) as source:
                audio = r.record(source)
                content = r.recognize_google(audio)
        elif file_type in ["mp4", "mov"]:
            clip = VideoFileClip(file_path)
            audio_path = "temp_audio.wav"
            clip.audio.write_audiofile(audio_path)
            r = sr.Recognizer()
            with sr.AudioFile(audio_path) as source:
                audio = r.record(source)
                content = r.recognize_google(audio)
            os.remove(audio_path)
        else:
            return file_type, "Unsupported file type."
    except Exception as e:
        return file_type, f"Error reading file: {str(e)}"

    return file_type, content

def store_document(file_path: str, file_type: str, content: str):
    """Store document metadata and content in database."""
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO documents (path, type, content) VALUES (%s, %s, %s) RETURNING id;",
        (file_path, file_type, content)
    )
    doc_id = cur.fetchone()[0]
    conn.commit()
    cur.close()
    conn.close()
    return doc_id

def answer_question(question: str, content: str) -> str:
    """Answer questions about document content using NLP."""
    try:
        result = qa_pipeline(question=question, context=content)
        return result["answer"]
    except Exception as e:
        return f"Error answering question: {str(e)}"

def read_document_agent(state):
    """Handle document reading and querying."""
    logger.info(f"Processing document command: {state.input}")
    command = state.input.lower().replace("teammate", "").strip()

    # Hardcoded file for now (enhance with Teams fetch later)
    if "excel" in command or "xlsx" in command:
        file_path = "data/sample.xlsx"
    elif "pdf" in command:
        file_path = "data/sample.pdf"
    else:
        state.context["response"] = "Please specify a document type (e.g., 'read the Excel file')."
        return state

    # Read and store document
    file_type, content = read_document(file_path)
    if "Error" in content or "Unsupported" in content:
        state.context["response"] = content
        return state
    
    doc_id = store_document(file_path, file_type, content)
    logger.info(f"Stored document ID {doc_id} with type {file_type}")

    # Answer question if present
    if "?" in command:
        question = command.split("read")[1].strip()
        answer = answer_question(question, content)
        state.context["response"] = f"Document ID {doc_id}: {answer}"
    else:
        state.context["response"] = f"Read document ID {doc_id}. Content: {content[:100]}..."

    return state

if __name__ == "__main__":
    state = {"input": "Teammate, what’s in row 5 of the Excel file?", "context": {}, "action_taken": True}
    result = read_document_agent(state)
    print(result)