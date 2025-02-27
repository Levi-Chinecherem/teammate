import os
import azure.cognitiveservices.speech as speechsdk
from msgraph.core import GraphClient
import requests
from pptx import Presentation
from utils.db import get_db_connection
from dotenv import load_dotenv
import logging
import time

# Load environment variables
load_dotenv()
client_id = os.getenv("CLIENT_ID")
client_secret = os.getenv("CLIENT_SECRET")
tenant_id = os.getenv("TENANT_ID")
speech_key = os.getenv("AZURE_SPEECH_KEY")
speech_region = os.getenv("AZURE_SPEECH_REGION")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("Presenter")

# Graph API client
def get_graph_client():
    auth_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    auth_data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": "https://graph.microsoft.com/.default"
    }
    response = requests.post(auth_url, data=auth_data)
    token = response.json().get("access_token")
    return GraphClient(api_version="v1.0", access_token=token)

# Speech synthesis
def speak(text):
    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=speech_region)
    synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config)
    result = synthesizer.speak_text_async(text).get()
    if result.reason != speechsdk.ResultReason.SynthesizingAudioCompleted:
        logger.error(f"Speech synthesis failed: {result.reason}")

def deliver_presentation(state):
    """Deliver a presentation with visibility verification."""
    logger.info(f"Processing presentation command: {state.input}")
    command = state.input.lower().replace("teammate", "").strip()

    if "present" not in command or "slides" not in command:
        state.context["response"] = "Please specify slides to present (e.g., 'present the Q1 slides')."
        return state

    # Extract slide file name (hardcoded for now)
    if "q1 slides" in command:
        slide_path = "data/q1_slides.pptx"
    else:
        state.context["response"] = "Slide file not recognized. Only 'Q1 slides' supported for now."
        return state

    # Load presentation
    try:
        prs = Presentation(slide_path)
    except Exception as e:
        state.context["response"] = f"Failed to load slides: {str(e)}"
        return state

    # Get latest meeting (enhance with context later)
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT id FROM meetings ORDER BY start_time DESC LIMIT 1;")
    meeting_id = cur.fetchone()
    if not meeting_id:
        state.context["response"] = "No recent meeting found to present in."
        cur.close()
        conn.close()
        return state
    meeting_id = meeting_id[0]
    cur.close()
    conn.close()

    # Join meeting and present (simplified with messages for now)
    client = get_graph_client()
    chat_id = state.context.get("chat_id", None)  # Placeholder; enhance with real meeting chat
    if not chat_id:
        # Simulate joining by sending to user's chat
        chat = client.users["user@example.com"].chats.post({"chatType": "oneOnOne"})
        chat_id = chat.id

    for slide in prs.slides:
        content = " ".join([shape.text for shape in slide.shapes if shape.has_text_frame and shape.text])
        if not content:
            continue

        # Simulate screen-sharing with message
        client.chats[chat_id].messages.post({"body": {"content": f"Sharing slide: {content}"}})
        speak("Can you see the slides?")
        time.sleep(2)  # Wait for response (manual for now)

        # Simulate visibility check (enhance with real response later)
        response = "yes"  # Hardcoded; replace with Teams message parsing
        while "no" in response.lower():
            logger.info("Resharing slide due to visibility issue")
            client.chats[chat_id].messages.post({"body": {"content": f"Resharing slide: {content}"}})
            speak("Can you see the slides now?")
            time.sleep(2)
            response = "yes"  # Simulate success

        speak(content)
        logger.info(f"Presented slide: {content}")
        time.sleep(1)  # Pause between slides

    state.context["response"] = f"Presented slides from {slide_path} in meeting ID {meeting_id}."
    return state

if __name__ == "__main__":
    state = {"input": "Teammate, present the Q1 slides", "context": {}, "action_taken": True}
    result = deliver_presentation(state)
    print(result)